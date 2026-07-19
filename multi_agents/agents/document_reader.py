"""Agente DocumentReader - Lectura de documentos en múltiples formatos"""

from pathlib import Path
from typing import Dict, Any, Optional, List, Union
import json

from tools.pdf_utils import extract_text_from_pdf, extract_metadata_from_pdf
from tools.odt_utils import (
    extract_text_from_odt,
    extract_text_from_ods,
    extract_text_from_odp,
    get_odf_metadata
)
from tools.vision_utils import analyze_image_with_llm, is_image_file


class DocumentReaderAgent:
    """
    Agente especializado en leer y extraer contenido de documentos.
    Soporta PDF, ODT, ODS, ODP, ODG, DOCX, XLSX, PPTX e imágenes.
    """
    
    def __init__(self, vision_client=None, vision_model: str = "local-vision-model"):
        """
        Inicializa el agente lector.
        
        Args:
            vision_client: Cliente para modelos de visión (opcional)
            vision_model: Nombre del modelo de visión
        """
        self.vision_client = vision_client
        self.vision_model = vision_model
    
    def read(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Lee un documento y extrae su contenido.
        
        Args:
            file_path: Ruta al archivo
        
        Returns:
            Diccionario con contenido y metadatos
        """
        path = Path(file_path)
        
        if not path.exists():
            return {"error": f"Archivo no encontrado: {file_path}"}
        
        extension = path.suffix.lower()
        
        # Dispatcher basado en extensión
        readers = {
            '.pdf': self._read_pdf,
            '.odt': self._read_odt,
            '.ods': self._read_ods,
            '.odp': self._read_odp,
            '.odg': self._read_odg,
            '.docx': self._read_docx,
            '.doc': self._read_doc,
            '.xlsx': self._read_xlsx,
            '.xls': self._read_xls,
            '.pptx': self._read_pptx,
            '.ppt': self._read_ppt,
            '.txt': self._read_txt,
            '.md': self._read_txt,
        }
        
        # Verificar si es imagen
        if is_image_file(path):
            return self._read_image(path)
        
        reader_func = readers.get(extension)
        
        if reader_func:
            return reader_func(path)
        else:
            return {
                "error": f"Formato no soportado: {extension}",
                "path": str(path),
                "supported_formats": list(readers.keys())
            }
    
    def _read_pdf(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo PDF"""
        text = extract_text_from_pdf(path)
        metadata = extract_metadata_from_pdf(path)
        
        return {
            "type": "pdf",
            "path": str(path),
            "content": text,
            "metadata": metadata,
            "pages": metadata.get('pages', 0),
            "success": True
        }
    
    def _read_odt(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo ODT"""
        text = extract_text_from_odt(path)
        metadata = get_odf_metadata(path)
        
        return {
            "type": "odt",
            "path": str(path),
            "content": text,
            "metadata": metadata,
            "success": True
        }
    
    def _read_ods(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo ODS"""
        text = extract_text_from_ods(path)
        metadata = get_odf_metadata(path)
        
        return {
            "type": "ods",
            "path": str(path),
            "content": text,
            "metadata": metadata,
            "success": True
        }
    
    def _read_odp(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo ODP"""
        text = extract_text_from_odp(path)
        metadata = get_odf_metadata(path)
        
        return {
            "type": "odp",
            "path": str(path),
            "content": text,
            "metadata": metadata,
            "success": True
        }
    
    def _read_odg(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo ODG (dibujo)"""
        # ODG es principalmente gráfico, usar visión si está disponible
        if self.vision_client:
            return self._read_image(path)
        else:
            metadata = get_odf_metadata(path)
            return {
                "type": "odg",
                "path": str(path),
                "content": "[ODG - Archivo gráfico. Use modelo de visión para análisis]",
                "metadata": metadata,
                "success": True,
                "requires_vision": True
            }
    
    def _read_docx(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo DOCX"""
        try:
            from docx import Document
            doc = Document(path)
            text = '\n'.join([para.text for para in doc.paragraphs])
            
            return {
                "type": "docx",
                "path": str(path),
                "content": text,
                "metadata": {" paragraphs": len(doc.paragraphs)},
                "success": True
            }
        except ImportError:
            return {
                "error": "python-docx no instalado. pip install python-docx",
                "path": str(path)
            }
    
    def _read_doc(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo DOC (antiguo formato Word)"""
        return {
            "error": "Formato DOC antiguo requiere conversión previa a DOCX o PDF",
            "path": str(path),
            "suggestion": "Convertir a PDF o DOCX primero"
        }
    
    def _read_xlsx(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo XLSX"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path)
            
            content = {}
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    rows.append(list(row))
                content[sheet_name] = rows
            
            return {
                "type": "xlsx",
                "path": str(path),
                "content": json.dumps(content, default=str),
                "metadata": {
                    "sheets": wb.sheetnames,
                    "active": wb.active.title
                },
                "success": True
            }
        except ImportError:
            return {
                "error": "openpyxl no instalado. pip install openpyxl",
                "path": str(path)
            }
    
    def _read_xls(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo XLS (antiguo formato Excel)"""
        return {
            "error": "Formato XLS antiguo requiere conversión previa",
            "path": str(path),
            "suggestion": "Convertir a XLSX o CSV primero"
        }
    
    def _read_pptx(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo PPTX"""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_text.append(f"Slide {i+1}: {' '.join(slide_text)}")
            
            return {
                "type": "pptx",
                "path": str(path),
                "content": '\n'.join(slides_text),
                "metadata": {"slides": len(prs.slides)},
                "success": True
            }
        except ImportError:
            return {
                "error": "python-pptx no instalado. pip install python-pptx",
                "path": str(path)
            }
    
    def _read_ppt(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo PPT (antiguo formato PowerPoint)"""
        return {
            "error": "Formato PPT antiguo requiere conversión previa",
            "path": str(path),
            "suggestion": "Convertir a PPTX o PDF primero"
        }
    
    def _read_txt(self, path: Path) -> Dict[str, Any]:
        """Lee un archivo de texto plano"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            return {
                "type": "text",
                "path": str(path),
                "content": text,
                "metadata": {"lines": text.count('\n') + 1},
                "success": True
            }
        except UnicodeDecodeError:
            with open(path, 'r', encoding='latin-1') as f:
                text = f.read()
            
            return {
                "type": "text",
                "path": str(path),
                "content": text,
                "metadata": {"lines": text.count('\n') + 1, "encoding": "latin-1"},
                "success": True
            }
    
    def _read_image(self, path: Path) -> Dict[str, Any]:
        """Lee una imagen usando el modelo de visión"""
        if not self.vision_client:
            return {
                "error": "No hay cliente de visión configurado",
                "path": str(path),
                "requires_vision": True
            }
        
        description = analyze_image_with_llm(
            self.vision_client,
            path,
            "Describe detalladamente el contenido de esta imagen, incluyendo cualquier texto visible.",
            self.vision_model
        )
        
        return {
            "type": "image",
            "path": str(path),
            "content": description,
            "metadata": {"format": path.suffix.lower()},
            "success": True
        }
    
    def read_multiple(self, file_paths: List[Union[str, Path]]) -> List[Dict[str, Any]]:
        """
        Lee múltiples documentos.
        
        Args:
            file_paths: Lista de rutas a archivos
        
        Returns:
            Lista de resultados de lectura
        """
        results = []
        for path in file_paths:
            result = self.read(path)
            results.append(result)
        return results
