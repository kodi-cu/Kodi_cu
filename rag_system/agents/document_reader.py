"""
Agente especializado en lectura y extracción de contenido de documentos.

Soporta múltiples formatos:
- PDF (pypdf, pdfplumber para tablas)
- ODT, ODS, ODP, ODG (OpenDocument)
- TXT, MD, CSV
- DOCX, XLSX, PPTX (Microsoft Office)
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass
import zipfile
import xml.etree.ElementTree as ET


@dataclass
class DocumentContent:
    """Contenido extraído de un documento."""
    path: str
    content: str
    metadata: Dict[str, Any]
    pages: Optional[List[str]] = None  # Para documentos paginados
    tables: Optional[List[Dict[str, Any]]] = None  # Tablas extraídas
    images_info: Optional[List[Dict[str, Any]]] = None  # Información de imágenes


class DocumentReaderAgent:
    """
    Agente para leer y extraer contenido de documentos.
    
    Soporta:
    - Extracción de texto plano
    - Detección de estructura (páginas, secciones)
    - Extracción de tablas
    - Metadatos del documento
    """
    
    def __init__(self):
        """Inicializa el agente lector."""
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Verifica dependencias instaladas."""
        self.dependencies = {
            'pypdf': False,
            'pdfplumber': False,
            'odfpy': False,
            'python_docx': False,
            'openpyxl': False,
            'python_pptx': False,
        }
        
        try:
            from pypdf import PdfReader
            self.dependencies['pypdf'] = True
        except ImportError:
            pass
        
        try:
            import pdfplumber
            self.dependencies['pdfplumber'] = True
        except ImportError:
            pass
        
        try:
            import odf.opendocument
            self.dependencies['odfpy'] = True
        except ImportError:
            pass
        
        try:
            import docx
            self.dependencies['python_docx'] = True
        except ImportError:
            pass
        
        try:
            import openpyxl
            self.dependencies['openpyxl'] = True
        except ImportError:
            pass
        
        try:
            import pptx
            self.dependencies['python_pptx'] = True
        except ImportError:
            pass
    
    def read_document(self, file_path: str) -> DocumentContent:
        """
        Lee un documento y extrae su contenido.
        
        Args:
            file_path: Ruta al archivo
            
        Returns:
            DocumentContent con el contenido extraído
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"Archivo no encontrado: {file_path}")
        
        ext = file_path.suffix.lower().lstrip('.')
        
        # Mapeo de extensiones a métodos
        readers = {
            'pdf': self._read_pdf,
            'odt': self._read_odt,
            'ods': self._read_ods,
            'odp': self._read_odp,
            'odg': self._read_odg,
            'txt': self._read_txt,
            'md': self._read_txt,
            'markdown': self._read_txt,
            'csv': self._read_csv,
            'docx': self._read_docx,
            'xlsx': self._read_xlsx,
            'pptx': self._read_pptx,
        }
        
        if ext not in readers:
            raise ValueError(f"Formato no soportado: {ext}")
        
        return readers[ext](str(file_path))
    
    def _read_pdf(self, file_path: str) -> DocumentContent:
        """Lee archivo PDF."""
        if not self.dependencies['pypdf']:
            raise ImportError("pypdf no está instalado. pip install pypdf")
        
        from pypdf import PdfReader
        
        reader = PdfReader(file_path)
        pages_content = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text.strip():
                pages_content.append(text)
        
        # Extraer metadatos
        metadata = {}
        if reader.metadata:
            metadata = {
                'title': reader.metadata.get('/Title', ''),
                'author': reader.metadata.get('/Author', ''),
                'subject': reader.metadata.get('/Subject', ''),
                'creator': reader.metadata.get('/Creator', ''),
                'producer': reader.metadata.get('/Producer', ''),
            }
        
        return DocumentContent(
            path=file_path,
            content='\n\n'.join(pages_content),
            metadata={
                **metadata,
                'source': file_path,
                'file_type': 'pdf',
                'total_pages': len(reader.pages),
            },
            pages=pages_content,
        )
    
    def _read_odt(self, file_path: str) -> DocumentContent:
        """Lee archivo ODT (OpenDocument Text)."""
        if not self.dependencies['odfpy']:
            # Fallback: leer como ZIP y extraer content.xml
            return self._read_opendocument_fallback(file_path, 'content.xml')
        
        import odf.opendocument
        import odf.text
        
        doc = odf.opendocument.load(file_path)
        content_parts = []
        
        # Recorrer todos los elementos de texto
        for element in doc.getElementsByType(odf.text.P):
            text = str(element)
            if text.strip():
                content_parts.append(text)
        
        return DocumentContent(
            path=file_path,
            content='\n'.join(content_parts),
            metadata={
                'source': file_path,
                'file_type': 'odt',
            },
        )
    
    def _read_ods(self, file_path: str) -> DocumentContent:
        """Lee archivo ODS (OpenDocument Spreadsheet)."""
        return self._read_opendocument_fallback(file_path, 'content.xml')
    
    def _read_odp(self, file_path: str) -> DocumentContent:
        """Lee archivo ODP (OpenDocument Presentation)."""
        return self._read_opendocument_fallback(file_path, 'content.xml')
    
    def _read_odg(self, file_path: str) -> DocumentContent:
        """Lee archivo ODG (OpenDocument Graphics)."""
        return self._read_opendocument_fallback(file_path, 'content.xml')
    
    def _read_opendocument_fallback(self, file_path: str, xml_file: str) -> DocumentContent:
        """
        Lectura fallback para archivos OpenDocument usando ZIP.
        
        Los archivos ODF son esencialmente archivos ZIP con XML dentro.
        """
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                if xml_file not in zip_file.namelist():
                    return DocumentContent(
                        path=file_path,
                        content="[Error: No se pudo extraer el contenido]",
                        metadata={
                            'source': file_path,
                            'file_type': 'opendocument',
                            'error': f'Archivo {xml_file} no encontrado',
                        },
                    )
                
                xml_content = zip_file.read(xml_file)
                
                # Parsear XML para extraer texto
                root = ET.fromstring(xml_content)
                text_parts = []
                
                # Buscar todos los elementos de texto
                namespaces = {
                    'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0',
                    'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
                    'draw': 'urn:oasis:names:tc:opendocument:xmlns:drawing:1.0',
                }
                
                # Extraer párrafos
                for elem in root.iter():
                    if 'paragraph' in elem.tag or 'h' in elem.tag:
                        if elem.text and elem.text.strip():
                            text_parts.append(elem.text.strip())
                
                # Contar imágenes
                images_info = []
                for name in zip_file.namelist():
                    if name.startswith('Pictures/'):
                        images_info.append({
                            'name': name,
                            'type': 'image',
                        })
                
                return DocumentContent(
                    path=file_path,
                    content='\n'.join(text_parts) if text_parts else "[Contenido XML crudo disponible]",
                    metadata={
                        'source': file_path,
                        'file_type': 'opendocument',
                        'xml_file': xml_file,
                    },
                    images_info=images_info if images_info else None,
                )
        except Exception as e:
            return DocumentContent(
                path=file_path,
                content=f"[Error leyendo archivo: {str(e)}]",
                metadata={
                    'source': file_path,
                    'file_type': 'opendocument',
                    'error': str(e),
                },
            )
    
    def _read_txt(self, file_path: str) -> DocumentContent:
        """Lee archivo de texto plano."""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        content = None
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            content = "[Error: No se pudo decodificar el archivo]"
        
        return DocumentContent(
            path=file_path,
            content=content,
            metadata={
                'source': file_path,
                'file_type': 'txt',
                'encoding': encoding,
            },
        )
    
    def _read_csv(self, file_path: str) -> DocumentContent:
        """Lee archivo CSV."""
        import csv
        
        rows = []
        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        
        # Convertir a formato legible
        content_lines = []
        for i, row in enumerate(rows):
            if i == 0:
                content_lines.append("Encabezados: " + ", ".join(row))
            else:
                content_lines.append("Fila " + str(i) + ": " + ", ".join(row))
        
        return DocumentContent(
            path=file_path,
            content='\n'.join(content_lines),
            metadata={
                'source': file_path,
                'file_type': 'csv',
                'total_rows': len(rows),
                'columns': len(rows[0]) if rows else 0,
            },
            tables=[{'headers': rows[0] if rows else [], 'data': rows[1:] if len(rows) > 1 else []}],
        )
    
    def _read_docx(self, file_path: str) -> DocumentContent:
        """Lee archivo DOCX."""
        if not self.dependencies['python_docx']:
            return DocumentContent(
                path=file_path,
                content="[python-docx no está instalado]",
                metadata={'source': file_path, 'file_type': 'docx'},
            )
        
        import docx
        
        doc = docx.Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        # Extraer tablas
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                tables.append({'data': table_data})
        
        return DocumentContent(
            path=file_path,
            content='\n\n'.join(paragraphs),
            metadata={
                'source': file_path,
                'file_type': 'docx',
                'total_paragraphs': len(paragraphs),
                'total_tables': len(tables),
            },
            tables=tables if tables else None,
        )
    
    def _read_xlsx(self, file_path: str) -> DocumentContent:
        """Lee archivo XLSX."""
        if not self.dependencies['openpyxl']:
            return DocumentContent(
                path=file_path,
                content="[openpyxl no está instalado]",
                metadata={'source': file_path, 'file_type': 'xlsx'},
            )
        
        import openpyxl
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets_content = []
        tables = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append([str(cell) if cell is not None else '' for cell in row])
            
            if rows:
                sheet_text = f"Hoja: {sheet_name}\n"
                for i, row in enumerate(rows):
                    if i == 0:
                        sheet_text += "Encabezados: " + ", ".join(row) + "\n"
                    else:
                        sheet_text += f"Fila {i}: " + ", ".join(row) + "\n"
                sheets_content.append(sheet_text)
                tables.append({'sheet': sheet_name, 'data': rows})
        
        return DocumentContent(
            path=file_path,
            content='\n\n'.join(sheets_content),
            metadata={
                'source': file_path,
                'file_type': 'xlsx',
                'total_sheets': len(wb.sheetnames),
                'sheet_names': wb.sheetnames,
            },
            tables=tables if tables else None,
        )
    
    def _read_pptx(self, file_path: str) -> DocumentContent:
        """Lee archivo PPTX."""
        if not self.dependencies['python_pptx']:
            return DocumentContent(
                path=file_path,
                content="[python-pptx no está instalado]",
                metadata={'source': file_path, 'file_type': 'pptx'},
            )
        
        import pptx
        
        prs = pptx.Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = f"Diapositiva {i}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            slides_content.append(slide_text)
        
        return DocumentContent(
            path=file_path,
            content='\n\n'.join(slides_content),
            metadata={
                'source': file_path,
                'file_type': 'pptx',
                'total_slides': len(prs.slides),
            },
            pages=slides_content,
        )
    
    def read_multiple(self, file_paths: List[str]) -> List[DocumentContent]:
        """
        Lee múltiples documentos.
        
        Args:
            file_paths: Lista de rutas a archivos
            
        Returns:
            Lista de DocumentContent
        """
        results = []
        for path in file_paths:
            try:
                content = self.read_document(path)
                results.append(content)
            except Exception as e:
                results.append(DocumentContent(
                    path=path,
                    content=f"[Error: {str(e)}]",
                    metadata={'source': path, 'error': str(e)},
                ))
        return results
    
    def to_context_string(self, contents: List[DocumentContent]) -> str:
        """
        Convierte múltiples contenidos en un string de contexto para LLM.
        
        Args:
            contents: Lista de DocumentContent
            
        Returns:
            String formateado para contexto
        """
        context_parts = []
        
        for i, content in enumerate(contents, 1):
            header = f"""
=== DOCUMENTO {i} ===
Archivo: {content.path}
Tipo: {content.metadata.get('file_type', 'desconocido')}
"""
            if 'total_pages' in content.metadata:
                header += f"Páginas: {content.metadata['total_pages']}\n"
            if 'total_slides' in content.metadata:
                header += f"Diapositivas: {content.metadata['total_slides']}\n"
            
            header += "=" * 40 + "\n\n"
            
            context_parts.append(header + content.content[:5000])  # Limitar longitud
        
        return "\n\n".join(context_parts)
